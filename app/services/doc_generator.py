"""
Document Generator — Creates professional Word (.docx) and PowerPoint (.pptx) files from AI content.
Supports full markdown: headings, bold/italic, code blocks, tables, bullet/numbered lists, blockquotes.
"""
import io
import re
import logging
from datetime import datetime

logger = logging.getLogger(__name__)

# ─── Color Constants ───
INDIGO = (99, 102, 241)
DARK_BG = (30, 32, 48)
WHITE = (255, 255, 255)
GRAY = (128, 128, 128)
LIGHT_GRAY = (160, 165, 180)
BODY_TEXT = (220, 222, 235)
CODE_BG = (245, 245, 250)
QUOTE_BAR = (99, 102, 241)
TABLE_HEADER_BG = (59, 130, 246)
TABLE_ALT_BG = (248, 250, 252)


def _clean_emoji_noise(content: str) -> str:
    """Remove UI button text that may leak into content."""
    # Remove common UI artifacts
    content = re.sub(r'📋\s*Copy\s*📄\s*Word\s*📊\s*PowerPoint\s*🔄\s*Retry\s*$', '', content, flags=re.MULTILINE)
    content = re.sub(r'📋\s*Copy\s*📄\s*Download Word\s*📊\s*Download PowerPoint\s*🔄\s*Retry\s*$', '', content, flags=re.MULTILINE)
    return content.strip()


# ═══════════════════════════════════════════════════════════════════
#  WORD DOCUMENT GENERATOR
# ═══════════════════════════════════════════════════════════════════

def generate_docx(content: str, title: str = "Document") -> io.BytesIO:
    """Generate a professional Word document from markdown content."""
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor, Cm, Emu
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.section import WD_ORIENT
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import parse_xml

    content = _clean_emoji_noise(content)
    doc = Document()

    # ── Page Setup ──
    section = doc.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.top_margin = Inches(0.8)
    section.bottom_margin = Inches(0.6)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)

    # ── Default Styles ──
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)
    style.font.color.rgb = RGBColor(40, 40, 40)
    style.paragraph_format.space_after = Pt(4)
    style.paragraph_format.line_spacing = 1.15

    # Heading styles
    for level in range(1, 5):
        try:
            h_style = doc.styles[f"Heading {level}"]
            h_style.font.name = "Calibri"
            h_style.font.color.rgb = RGBColor(30, 41, 59)
            if level == 1:
                h_style.font.size = Pt(20)
                h_style.paragraph_format.space_before = Pt(18)
                h_style.paragraph_format.space_after = Pt(8)
            elif level == 2:
                h_style.font.size = Pt(16)
                h_style.paragraph_format.space_before = Pt(14)
                h_style.paragraph_format.space_after = Pt(6)
            elif level == 3:
                h_style.font.size = Pt(13)
                h_style.paragraph_format.space_before = Pt(10)
                h_style.paragraph_format.space_after = Pt(4)
            else:
                h_style.font.size = Pt(11)
                h_style.paragraph_format.space_before = Pt(8)
                h_style.paragraph_format.space_after = Pt(4)
        except KeyError:
            pass

    # ── Title Block ──
    # Accent bar above title
    bar_para = doc.add_paragraph()
    bar_para.paragraph_format.space_after = Pt(2)
    bar_run = bar_para.add_run("━" * 60)
    bar_run.font.color.rgb = RGBColor(*INDIGO)
    bar_run.font.size = Pt(6)

    # Title
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
    title_para.paragraph_format.space_after = Pt(4)
    run = title_para.add_run(title)
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = RGBColor(15, 23, 42)
    run.font.name = "Calibri"

    # Subtitle with date & org
    sub_para = doc.add_paragraph()
    sub_para.paragraph_format.space_after = Pt(2)
    run = sub_para.add_run(f"General Department of Taxation, Cambodia")
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(*GRAY)

    date_para = doc.add_paragraph()
    date_para.paragraph_format.space_after = Pt(6)
    run = date_para.add_run(f"Generated on {datetime.now().strftime('%B %d, %Y at %I:%M %p')}")
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(*GRAY)
    run.font.italic = True

    # Bottom accent bar
    bar2 = doc.add_paragraph()
    bar2.paragraph_format.space_after = Pt(12)
    bar2_run = bar2.add_run("━" * 60)
    bar2_run.font.color.rgb = RGBColor(*INDIGO)
    bar2_run.font.size = Pt(6)

    # ── Parse & Render Content ──
    lines = content.split("\n")
    in_code_block = False
    code_lines = []
    code_lang = ""
    in_table = False
    table_rows = []

    for line in lines:
        # Code block handling
        if line.strip().startswith("```"):
            if in_code_block:
                code_text = "\n".join(code_lines)
                _add_code_block(doc, code_text, code_lang)
                in_code_block = False
                code_lines = []
                code_lang = ""
            else:
                in_code_block = True
                code_lang = line.strip().replace("```", "").strip()
            continue

        if in_code_block:
            code_lines.append(line)
            continue

        # Table handling
        if "|" in line and line.strip().startswith("|"):
            cells = [c.strip() for c in line.split("|") if c.strip()]
            if all(set(c) <= set("- :") for c in cells):
                continue  # Skip separator row
            table_rows.append(cells)
            in_table = True
            continue
        elif in_table:
            if table_rows:
                _add_docx_table(doc, table_rows)
            table_rows = []
            in_table = False

        stripped = line.strip()

        # Headings
        if stripped.startswith("#### "):
            _add_heading_with_emoji(doc, stripped[5:], 4)
        elif stripped.startswith("### "):
            _add_heading_with_emoji(doc, stripped[4:], 3)
        elif stripped.startswith("## "):
            _add_heading_with_emoji(doc, stripped[3:], 2)
        elif stripped.startswith("# "):
            _add_heading_with_emoji(doc, stripped[2:], 1)
        # Horizontal rule
        elif stripped.startswith("---") or stripped.startswith("***"):
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(8)
            p.paragraph_format.space_after = Pt(8)
            run = p.add_run("─" * 50)
            run.font.color.rgb = RGBColor(200, 200, 210)
            run.font.size = Pt(8)
        # Bullet list
        elif stripped.startswith("- ") or stripped.startswith("* "):
            text = stripped[2:]
            p = doc.add_paragraph(style="List Bullet")
            p.paragraph_format.space_after = Pt(2)
            _add_formatted_text(p, text)
        # Sub-bullet (indented)
        elif re.match(r"^\s{2,}[-*]\s", line):
            text = line.strip()[2:]
            p = doc.add_paragraph(style="List Bullet 2") if "List Bullet 2" in [s.name for s in doc.styles] else doc.add_paragraph(style="List Bullet")
            p.paragraph_format.left_indent = Inches(0.5)
            p.paragraph_format.space_after = Pt(2)
            _add_formatted_text(p, text)
        # Numbered list
        elif re.match(r"^\d+\.\s", stripped):
            text = re.sub(r"^\d+\.\s", "", stripped)
            p = doc.add_paragraph(style="List Number")
            p.paragraph_format.space_after = Pt(2)
            _add_formatted_text(p, text)
        # Blockquote
        elif stripped.startswith("> "):
            text = stripped[2:]
            _add_blockquote(doc, text)
        # Checkbox items
        elif stripped.startswith("- [ ] ") or stripped.startswith("- [x] "):
            checked = stripped.startswith("- [x] ")
            text = stripped[6:]
            marker = "☑" if checked else "☐"
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Inches(0.3)
            p.paragraph_format.space_after = Pt(2)
            run = p.add_run(f"{marker} ")
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(*INDIGO) if checked else RGBColor(*GRAY)
            _add_formatted_text(p, text)
        # Empty line — add spacing
        elif not stripped:
            p = doc.add_paragraph()
            p.paragraph_format.space_after = Pt(2)
        # Normal paragraph
        else:
            p = doc.add_paragraph()
            p.paragraph_format.space_after = Pt(4)
            _add_formatted_text(p, stripped)

    # Flush remaining
    if in_code_block and code_lines:
        _add_code_block(doc, "\n".join(code_lines), code_lang)
    if table_rows:
        _add_docx_table(doc, table_rows)

    # ── Footer ──
    doc.add_paragraph("")
    bar3 = doc.add_paragraph()
    bar3_run = bar3.add_run("━" * 60)
    bar3_run.font.color.rgb = RGBColor(*INDIGO)
    bar3_run.font.size = Pt(6)

    footer = doc.add_paragraph()
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer.paragraph_format.space_before = Pt(4)
    run = footer.add_run("Generated by AI Personal Assistant — aia.rikreay24.com")
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(160, 160, 160)
    run.font.italic = True

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def _add_heading_with_emoji(doc, text, level):
    """Add heading, preserving emoji characters."""
    from docx.shared import Pt, RGBColor
    h = doc.add_heading(level=level)
    # Clean markdown bold from heading text
    clean_text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    clean_text = re.sub(r'\*(.+?)\*', r'\1', clean_text)
    run = h.add_run(clean_text)
    run.font.name = "Calibri"
    if level == 1:
        run.font.color.rgb = RGBColor(15, 23, 42)
    elif level == 2:
        run.font.color.rgb = RGBColor(30, 58, 138)
    else:
        run.font.color.rgb = RGBColor(51, 65, 85)


def _add_code_block(doc, code_text, lang=""):
    """Add a formatted code block with background shading."""
    from docx.shared import Inches, Pt, RGBColor
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import parse_xml

    if lang:
        label = doc.add_paragraph()
        label.paragraph_format.space_after = Pt(0)
        label.paragraph_format.space_before = Pt(6)
        run = label.add_run(f"  {lang}")
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(*INDIGO)
        run.font.bold = True

    p = doc.add_paragraph()
    p.paragraph_format.left_indent = Inches(0.2)
    p.paragraph_format.right_indent = Inches(0.2)
    p.paragraph_format.space_before = Pt(4)
    p.paragraph_format.space_after = Pt(8)

    # Add background shading
    try:
        shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="F1F5F9" w:val="clear"/>')
        p._element.get_or_add_pPr().append(shading)
    except Exception:
        pass

    run = p.add_run(code_text)
    run.font.name = "Consolas"
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(30, 41, 59)


def _add_blockquote(doc, text):
    """Add a styled blockquote."""
    from docx.shared import Inches, Pt, RGBColor
    from docx.oxml.ns import nsdecls
    from docx.oxml import parse_xml

    p = doc.add_paragraph()
    p.paragraph_format.left_indent = Inches(0.4)
    p.paragraph_format.space_before = Pt(4)
    p.paragraph_format.space_after = Pt(4)

    # Add left border shading
    try:
        shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="EEF2FF" w:val="clear"/>')
        p._element.get_or_add_pPr().append(shading)
    except Exception:
        pass

    bar_run = p.add_run("│ ")
    bar_run.font.color.rgb = RGBColor(*INDIGO)
    bar_run.font.size = Pt(12)

    _add_formatted_text(p, text, italic=True, color=RGBColor(71, 85, 105))


def _add_docx_table(doc, rows):
    """Add a professionally styled table."""
    from docx.shared import Pt, RGBColor, Inches, Cm
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import parse_xml

    if not rows:
        return

    cols = max(len(r) for r in rows)
    table = doc.add_table(rows=len(rows), cols=cols)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER

    # Try to use a clean style
    try:
        table.style = "Table Grid"
    except Exception:
        pass

    for i, row_data in enumerate(rows):
        for j, cell_text in enumerate(row_data):
            if j < cols:
                cell = table.cell(i, j)
                cell.text = ""

                # Clean markdown from cell text
                clean = re.sub(r'\*\*(.+?)\*\*', r'\1', cell_text)
                clean = re.sub(r'\*(.+?)\*', r'\1', clean)

                p = cell.paragraphs[0]
                p.paragraph_format.space_before = Pt(3)
                p.paragraph_format.space_after = Pt(3)

                run = p.add_run(clean)
                run.font.size = Pt(10)
                run.font.name = "Calibri"

                if i == 0:
                    # Header row — bold white on blue
                    run.bold = True
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    try:
                        shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="3B5998" w:val="clear"/>')
                        cell._element.get_or_add_tcPr().append(shading)
                    except Exception:
                        pass
                else:
                    run.font.color.rgb = RGBColor(30, 41, 59)
                    # Alternating row colors
                    if i % 2 == 0:
                        try:
                            shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="F8FAFC" w:val="clear"/>')
                            cell._element.get_or_add_tcPr().append(shading)
                        except Exception:
                            pass

    doc.add_paragraph("")  # Spacer after table


def _add_formatted_text(paragraph, text, italic=False, color=None):
    """Add text with bold/italic/code markdown formatting to a paragraph."""
    from docx.shared import Pt, RGBColor

    default_color = color or RGBColor(40, 40, 40)

    # Split by formatting markers
    parts = re.split(r"(\*\*\*.+?\*\*\*|\*\*.+?\*\*|\*.+?\*|`.+?`)", text)
    for part in parts:
        if not part:
            continue
        if part.startswith("***") and part.endswith("***"):
            run = paragraph.add_run(part[3:-3])
            run.bold = True
            run.italic = True
            run.font.color.rgb = default_color
        elif part.startswith("**") and part.endswith("**"):
            run = paragraph.add_run(part[2:-2])
            run.bold = True
            run.font.color.rgb = default_color
        elif part.startswith("*") and part.endswith("*") and len(part) > 2:
            run = paragraph.add_run(part[1:-1])
            run.italic = True
            run.font.color.rgb = default_color
        elif part.startswith("`") and part.endswith("`"):
            run = paragraph.add_run(part[1:-1])
            run.font.name = "Consolas"
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(*INDIGO)
        else:
            run = paragraph.add_run(part)
            run.font.color.rgb = default_color
            if italic:
                run.italic = True


# ═══════════════════════════════════════════════════════════════════
#  POWERPOINT GENERATOR
# ═══════════════════════════════════════════════════════════════════

def generate_pptx(content: str, title: str = "Presentation") -> io.BytesIO:
    """Generate a professional PowerPoint presentation from markdown content."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    content = _clean_emoji_noise(content)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color scheme
    primary = RGBColor(*INDIGO)
    dark = RGBColor(*DARK_BG)
    white = RGBColor(*WHITE)
    gray = RGBColor(*LIGHT_GRAY)
    body_clr = RGBColor(*BODY_TEXT)
    accent2 = RGBColor(59, 130, 246)  # Blue

    # ── Title Slide ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _set_slide_bg(slide, dark)

    # Decorative accent — top bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = primary
    bar.line.fill.background()

    # Side accent bar
    side = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.06), Inches(7.5))
    side.fill.solid()
    side.fill.fore_color.rgb = primary
    side.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.333), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER

    # Accent line under title
    line = slide.shapes.add_shape(1, Inches(5), Inches(4.0), Inches(3.333), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = primary
    line.line.fill.background()

    # Subtitle — org
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(10.333), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "General Department of Taxation, Cambodia"
    p2.font.size = Pt(18)
    p2.font.color.rgb = gray
    p2.alignment = PP_ALIGN.CENTER

    # Date
    txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.1), Inches(10.333), Inches(0.6))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = datetime.now().strftime('%B %d, %Y')
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(100, 105, 120)
    p3.alignment = PP_ALIGN.CENTER

    # ── Content Slides ──
    sections = _parse_sections(content)

    for si, section in enumerate(sections):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_slide_bg(slide, dark)

        # Top accent bar
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = primary
        bar.line.fill.background()

        # Side accent
        side = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.06), Inches(7.5))
        side.fill.solid()
        side.fill.fore_color.rgb = primary
        side.line.fill.background()

        # Slide number
        num_box = slide.shapes.add_textbox(Inches(12.2), Inches(7.0), Inches(1), Inches(0.4))
        num_tf = num_box.text_frame
        num_p = num_tf.paragraphs[0]
        num_p.text = f"{si + 1}"
        num_p.font.size = Pt(11)
        num_p.font.color.rgb = RGBColor(80, 85, 100)
        num_p.alignment = PP_ALIGN.RIGHT

        heading = section.get("heading", "")
        body_lines = section.get("body", [])
        table_data = section.get("table", None)

        # Clean emoji from heading for cleaner look (keep them)
        heading_clean = re.sub(r'\*\*(.+?)\*\*', r'\1', heading)

        # Section heading
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = heading_clean
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = white

        # Heading underline
        h_line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.35), Inches(3), Inches(0.03))
        h_line.fill.solid()
        h_line.fill.fore_color.rgb = primary
        h_line.line.fill.background()

        # Body content
        if body_lines:
            txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.733), Inches(5.2))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True

            for i, line in enumerate(body_lines):
                if i == 0:
                    p = tf2.paragraphs[0]
                else:
                    p = tf2.add_paragraph()

                cleaned = line.strip()
                is_bullet = cleaned.startswith("- ") or cleaned.startswith("* ")
                is_numbered = bool(re.match(r"^\d+\.\s", cleaned))
                is_checkbox = cleaned.startswith("- [ ] ") or cleaned.startswith("- [x] ")

                if is_checkbox:
                    checked = cleaned.startswith("- [x] ")
                    cleaned = cleaned[6:]
                    marker = "☑ " if checked else "☐ "
                    p.text = marker + cleaned
                    p.level = 0
                elif is_bullet:
                    cleaned = cleaned[2:]
                    p.text = "▸  " + cleaned
                    p.level = 0
                elif is_numbered:
                    num_match = re.match(r"^(\d+)\.\s", cleaned)
                    num = num_match.group(1) if num_match else ""
                    cleaned = re.sub(r"^\d+\.\s", "", cleaned)
                    p.text = f"{num}.  " + cleaned
                    p.level = 0
                else:
                    p.text = cleaned

                # Remove markdown formatting from display text
                p.text = re.sub(r"\*\*\*(.+?)\*\*\*", r"\1", p.text)
                p.text = re.sub(r"\*\*(.+?)\*\*", r"\1", p.text)
                p.text = re.sub(r"\*(.+?)\*", r"\1", p.text)
                p.text = re.sub(r"`(.+?)`", r"\1", p.text)

                p.font.size = Pt(17)
                p.font.color.rgb = body_clr
                p.space_after = Pt(6)

                # Re-render with bold formatting
                if "**" in line:
                    original_text = p.text
                    p.clear()
                    # Get prefix for bullets
                    prefix = ""
                    if is_bullet:
                        prefix = "▸  "
                    elif is_numbered:
                        num_match = re.match(r"^(\d+)\.\s", line.strip())
                        prefix = f"{num_match.group(1)}.  " if num_match else ""

                    if prefix:
                        run = p.add_run()
                        run.text = prefix
                        run.font.size = Pt(17)
                        run.font.color.rgb = primary

                    # Process the cleaned text for bold
                    text_to_parse = cleaned if (is_bullet or is_numbered) else line.strip()
                    text_to_parse = re.sub(r"\*\*\*(.+?)\*\*\*", r"**\1**", text_to_parse)
                    parts = re.split(r"(\*\*.+?\*\*)", text_to_parse)
                    for part in parts:
                        if part.startswith("**") and part.endswith("**"):
                            run = p.add_run()
                            run.text = part[2:-2]
                            run.font.bold = True
                            run.font.size = Pt(17)
                            run.font.color.rgb = white
                        elif part:
                            run = p.add_run()
                            run.text = part
                            run.font.size = Pt(17)
                            run.font.color.rgb = body_clr

        # If section has a table, add it as a separate slide
        if table_data and len(table_data) > 0:
            _add_pptx_table_slide(prs, heading_clean, table_data, dark, primary, white, gray, body_clr)

    # ── Thank You Slide ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, dark)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = primary
    bar.line.fill.background()

    side = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.06), Inches(7.5))
    side.fill.solid()
    side.fill.fore_color.rgb = primary
    side.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER

    line = slide.shapes.add_shape(1, Inches(5), Inches(4.0), Inches(3.333), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = primary
    line.line.fill.background()

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.4), Inches(11.333), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "AI Personal Assistant — aia.rikreay24.com"
    p2.font.size = Pt(14)
    p2.font.color.rgb = gray
    p2.alignment = PP_ALIGN.CENTER

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _add_pptx_table_slide(prs, heading, table_data, dark, primary, white, gray, body_clr):
    """Add a slide with a formatted table."""
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, dark)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = primary
    bar.line.fill.background()

    side = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.06), Inches(7.5))
    side.fill.solid()
    side.fill.fore_color.rgb = primary
    side.line.fill.background()

    # Heading
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = heading
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = white

    # Table
    cols = max(len(r) for r in table_data)
    rows_count = len(table_data)
    col_width = min(Inches(11.4 / cols), Inches(4))

    table = slide.shapes.add_table(
        rows_count, cols,
        Inches(0.8), Inches(1.5),
        Inches(min(11.4, col_width * cols / 914400 * cols)), Inches(min(5.0, rows_count * 0.5))
    ).table

    for i, row_data in enumerate(table_data):
        for j, cell_text in enumerate(row_data):
            if j < cols:
                cell = table.cell(i, j)
                # Clean markdown
                clean = re.sub(r'\*\*(.+?)\*\*', r'\1', cell_text)
                clean = re.sub(r'\*(.+?)\*', r'\1', clean)

                p = cell.text_frame.paragraphs[0]
                p.text = clean
                p.font.size = Pt(14)

                if i == 0:
                    p.font.bold = True
                    p.font.color.rgb = white
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = primary
                else:
                    p.font.color.rgb = body_clr
                    if i % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(35, 38, 55)


def _set_slide_bg(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _parse_sections(content: str) -> list:
    """Parse markdown content into slide sections (split by headings).
    Also extracts tables as separate data."""
    sections = []
    current = None
    lines = content.split("\n")
    in_table = False
    table_rows = []

    for line in lines:
        stripped = line.strip()

        # Heading detection
        if stripped.startswith("## ") or stripped.startswith("### ") or stripped.startswith("# "):
            # Flush table to current section
            if in_table and table_rows and current:
                current["table"] = table_rows
                table_rows = []
                in_table = False

            if current:
                sections.append(current)
            heading = re.sub(r"^#{1,4}\s*", "", stripped)
            current = {"heading": heading, "body": [], "table": None}

        # Table row
        elif "|" in stripped and stripped.startswith("|"):
            cells = [c.strip() for c in stripped.split("|") if c.strip()]
            if all(set(c) <= set("- :") for c in cells):
                continue  # separator
            table_rows.append(cells)
            in_table = True

        elif in_table:
            # End of table
            if current:
                current["table"] = table_rows
            table_rows = []
            in_table = False
            # Process this line normally
            if current is not None:
                if stripped and not stripped.startswith("---"):
                    current["body"].append(stripped)
            elif stripped:
                current = {"heading": "Overview", "body": [stripped], "table": None}

        elif current is not None:
            if stripped and not stripped.startswith("---") and not stripped.startswith("```"):
                current["body"].append(stripped)
        else:
            if stripped and not stripped.startswith("---") and not stripped.startswith("```"):
                current = {"heading": "Overview", "body": [stripped], "table": None}

    # Flush remaining
    if in_table and table_rows and current:
        current["table"] = table_rows
    if current:
        sections.append(current)

    if not sections:
        sections = [{"heading": "Content", "body": content.split("\n"), "table": None}]

    return sections
